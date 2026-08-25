"""Tests for src.quantlib.valuation.artifact.

Structured around the six things the artifact module promises:

1. The input hash is stable under dict-key reordering and int/float/``-0.0``
   representation noise, and changes whenever any single input number changes
   (:class:`TestInputHashStability`).
2. The module never reads a clock -- checked both by grepping its own source
   (no ``datetime.now``/``time.time``) and by proving omission/naivety of
   ``generated_at`` raises rather than silently defaulting
   (:class:`TestNoHiddenClock`).
3. Every :class:`~src.quantlib.valuation.contracts.Assumption` in a result is
   captured with its ``basis``, and models that take no ``Assumption``
   objects (comps, three-statement) honestly report an empty tuple rather
   than fabricating one (:class:`TestAssumptionCapture`).
4. Two artifacts diff to the exact changed field and its numeric delta
   (:class:`TestDiff`).
5. Every xlsx export round-trips through ``openpyxl`` back to the values on
   the result object -- not just "the file exists" -- and the excluded/
   omitted/non-converged material is actually present in the workbook
   (:class:`TestDCFWorkbookExport`, :class:`TestCompsWorkbookExport`,
   :class:`TestThreeStatementWorkbookExport`).
6. The pptx summary carries the same input hash, assumptions and gap notes
   (:class:`TestSummaryDeckExport`).
"""

from __future__ import annotations

import ast
import inspect
from datetime import datetime, timedelta, timezone
from pathlib import Path

import openpyxl
import pandas as pd
import pytest
from pptx import Presentation

from src.quantlib.valuation import artifact as artifact_module
from src.quantlib.valuation.artifact import (
    AssumptionRecord,
    ModelArtifact,
    build_comps_artifact,
    build_dcf_artifact,
    build_three_statement_artifact,
    compute_input_hash,
    diff_artifacts,
    export_comps_workbook,
    export_dcf_workbook,
    export_summary_deck,
    export_three_statement_workbook,
)
from src.quantlib.valuation.comps import FlowMetricPeriods, PeerCompany, TargetCompany, run_comps
from src.quantlib.valuation.contracts import Assumption
from src.quantlib.valuation.dcf import run_dcf, sensitivity_grid
from src.quantlib.valuation.threestatement import project_three_statement

UTC = timezone.utc
FIXED_TIME = datetime(2026, 1, 15, 12, 0, 0, tzinfo=UTC)


# ---------------------------------------------------------------------------
# Shared fixtures (self-contained -- not imported from the sibling test files)
# ---------------------------------------------------------------------------


def _dcf_inputs() -> dict:
    """A small, valid ``run_dcf`` inputs mapping (target capital structure)."""
    return {
        "risk_free_rate": 0.04,
        "beta": 1.2,
        "equity_risk_premium": 0.05,
        "size_premium": 0.01,
        "country_risk_premium": 0.0,
        "pretax_cost_of_debt": 0.06,
        "tax_rate": 0.25,
        "ebit": [100.0, 105.0, 110.0],
        "depreciation_amortization": [20.0, 21.0, 22.0],
        "capex": [30.0, 28.0, 26.0],
        "delta_nwc": [5.0, 4.0, 3.0],
        "terminal_growth": Assumption(
            name="terminal_growth",
            value=0.03,
            basis="Long-run nominal GDP proxy for a mature-market perpetuity",
            source="analyst view",
        ),
        "exit_multiple": Assumption(
            name="exit_multiple",
            value=8.0,
            basis="Median EV/EBITDA of the comparable peer set",
        ),
        "total_debt": 200.0,
        "cash_and_equivalents": 50.0,
        "minority_interest": 10.0,
        "preferred_equity": 15.0,
        "associate_investments": 5.0,
        "diluted_shares": 100.0,
        "target_equity_weight": 0.7,
        "target_debt_weight": 0.3,
    }


def _run_dcf_default(inputs: dict | None = None):
    return run_dcf(
        inputs if inputs is not None else _dcf_inputs(),
        capital_structure_basis="target",
        discounting_convention="mid_year",
        terminal_value_method="perpetuity_growth",
    )


def _build_dcf_artifact(inputs: dict | None = None, *, generated_at: datetime = FIXED_TIME):
    inputs = inputs if inputs is not None else _dcf_inputs()
    result = _run_dcf_default(inputs)
    return build_dcf_artifact(
        inputs=inputs,
        result=result,
        generated_at=generated_at,
        capital_structure_basis="target",
        discounting_convention="mid_year",
        terminal_value_method="perpetuity_growth",
    )


def _flow(last_full: float, *, ytd: float, prior_ytd: float, month: int = 12) -> FlowMetricPeriods:
    return FlowMetricPeriods(
        fiscal_year_end_month=month,
        last_full_fiscal_year=last_full,
        current_year_to_date=ytd,
        prior_year_to_date=prior_ytd,
    )


def _comps_peers(*, include_excluded_peer: bool = True) -> list[PeerCompany]:
    peers = [
        PeerCompany(
            name="PEER_A",
            market_cap=1000.0,
            total_debt=200.0,
            cash_and_equivalents=100.0,
            ebitda=_flow(150.0, ytd=80.0, prior_ytd=70.0),
            ebit=_flow(100.0, ytd=55.0, prior_ytd=50.0),
            revenue=_flow(900.0, ytd=480.0, prior_ytd=420.0),
            diluted_eps=_flow(5.0, ytd=2.6, prior_ytd=2.2),
            price_per_share=54.0,
            book_value_of_equity=500.0,
            eps_basis="gaap",
            minority_interest=None,
            preferred_stock=None,
            investments_in_associates=None,
        ),
        PeerCompany(
            name="PEER_B",
            market_cap=1200.0,
            total_debt=250.0,
            cash_and_equivalents=120.0,
            ebitda=_flow(160.0, ytd=85.0, prior_ytd=75.0),
            ebit=_flow(110.0, ytd=58.0, prior_ytd=52.0),
            revenue=_flow(950.0, ytd=500.0, prior_ytd=440.0),
            diluted_eps=_flow(5.5, ytd=2.8, prior_ytd=2.4),
            price_per_share=60.0,
            book_value_of_equity=550.0,
            eps_basis="gaap",
        ),
    ]
    if include_excluded_peer:
        # Negative LTM EBITDA -> excluded from ev_ebitda's distribution.
        peers.append(
            PeerCompany(
                name="PEER_LOSSMAKER",
                market_cap=300.0,
                total_debt=400.0,
                cash_and_equivalents=20.0,
                ebitda=_flow(-50.0, ytd=-30.0, prior_ytd=-10.0),
                ebit=_flow(-60.0, ytd=-35.0, prior_ytd=-15.0),
                revenue=_flow(500.0, ytd=260.0, prior_ytd=240.0),
                diluted_eps=_flow(-1.0, ytd=-0.5, prior_ytd=-0.2),
                price_per_share=8.0,
                book_value_of_equity=90.0,
                eps_basis="gaap",
            )
        )
    return peers


def _comps_target(*, omit_components: bool = True) -> TargetCompany:
    kwargs = dict(
        name="TARGET",
        total_debt=150.0,
        cash_and_equivalents=60.0,
        ebitda=_flow(140.0, ytd=75.0, prior_ytd=65.0),
        ebit=_flow(95.0, ytd=52.0, prior_ytd=45.0),
        revenue=_flow(850.0, ytd=460.0, prior_ytd=400.0),
        diluted_eps=_flow(4.5, ytd=2.4, prior_ytd=2.0),
        diluted_shares_outstanding=100.0,
        book_value_of_equity=420.0,
        eps_basis="gaap",
    )
    if not omit_components:
        kwargs.update(minority_interest=0.0, preferred_stock=0.0, investments_in_associates=0.0)
    return TargetCompany(**kwargs)


def _build_comps_artifact(*, generated_at: datetime = FIXED_TIME):
    target = _comps_target()
    peers = _comps_peers()
    result = run_comps(target, peers, calendarisation_policy="ltm")
    art = build_comps_artifact(
        target=target,
        peers=peers,
        calendarisation_policy="ltm",
        result=result,
        generated_at=generated_at,
    )
    return target, peers, result, art


def _three_statement_opening() -> dict:
    return {
        "revenue": 1_000_000.0,
        "cash": 200_000.0,
        "net_working_capital": 150_000.0,
        "ppe": 800_000.0,
        "revolver_balance": 300_000.0,
        "paid_in_capital": 500_000.0,
        "retained_earnings": 350_000.0,
    }


def _three_statement_drivers() -> dict:
    # minimum_cash deliberately high in period 1 to force a revolver draw --
    # exercises the "declared plug" path the artifact must surface.
    return {
        "revenue_growth": [0.10, 0.08],
        "gross_margin": [0.40, 0.42],
        "opex_pct_revenue": [0.20, 0.19],
        "capex_pct_revenue": [0.05, 0.05],
        "nwc_pct_revenue": [0.15, 0.15],
        "tax_rate": [0.25, 0.25],
        "dividend_payout_ratio": [0.30, 0.30],
        "depreciation_amortization": [60_000.0, 65_000.0],
        "interest_rate": [0.06, 0.06],
        "minimum_cash": [400_000.0, 100_000.0],
    }


def _build_three_statement_artifact(*, generated_at: datetime = FIXED_TIME):
    opening = _three_statement_opening()
    drivers = _three_statement_drivers()
    result = project_three_statement(opening, drivers)
    art = build_three_statement_artifact(
        opening=opening, drivers=drivers, result=result, generated_at=generated_at
    )
    return opening, drivers, result, art


# ---------------------------------------------------------------------------
# 1. Input hash stability (both sides required by the task)
# ---------------------------------------------------------------------------


class TestInputHashStability:
    def test_same_payload_hashes_identically_twice(self):
        payload = {"a": 1.0, "b": {"c": 2.0, "d": [1, 2, 3]}}
        assert compute_input_hash(payload) == compute_input_hash(payload)

    def test_hash_independent_of_dict_key_order(self):
        forward = {"total_debt": 200.0, "cash": 50.0, "tax_rate": 0.25}
        reversed_ = {"tax_rate": 0.25, "cash": 50.0, "total_debt": 200.0}
        assert compute_input_hash(forward) == compute_input_hash(reversed_)

    def test_hash_independent_of_nested_dict_key_order(self):
        forward = {"outer": {"x": 1.0, "y": 2.0}, "z": 3.0}
        reordered = {"z": 3.0, "outer": {"y": 2.0, "x": 1.0}}
        assert compute_input_hash(forward) == compute_input_hash(reordered)

    def test_hash_treats_int_and_float_as_the_same_value(self):
        assert compute_input_hash({"shares": 100}) == compute_input_hash({"shares": 100.0})

    def test_hash_collapses_negative_zero(self):
        assert compute_input_hash({"x": -0.0}) == compute_input_hash({"x": 0.0})

    @pytest.mark.parametrize(
        ("field", "changed_value"),
        [
            ("total_debt", 201.0),
            ("cash_and_equivalents", 51.0),
            ("tax_rate", 0.26),
            ("risk_free_rate", 0.041),
            ("diluted_shares", 101.0),
        ],
    )
    def test_hash_changes_when_a_single_scalar_input_changes(self, field, changed_value):
        baseline = _dcf_inputs()
        changed = _dcf_inputs()
        changed[field] = changed_value
        assert compute_input_hash(baseline) != compute_input_hash(changed)

    def test_hash_changes_when_a_sequence_entry_changes(self):
        baseline = _dcf_inputs()
        changed = _dcf_inputs()
        changed["ebit"] = [100.0, 105.0, 111.0]  # last year bumped by 1.0
        assert compute_input_hash(baseline) != compute_input_hash(changed)

    def test_hash_changes_when_an_assumption_value_changes(self):
        baseline = _dcf_inputs()
        changed = _dcf_inputs()
        changed["terminal_growth"] = Assumption(
            name="terminal_growth", value=0.031, basis=baseline["terminal_growth"].basis
        )
        assert compute_input_hash(baseline) != compute_input_hash(changed)

    def test_hash_changes_when_only_an_assumption_basis_changes(self):
        baseline = _dcf_inputs()
        changed = _dcf_inputs()
        changed["terminal_growth"] = Assumption(
            name="terminal_growth", value=0.03, basis="A completely different justification"
        )
        assert compute_input_hash(baseline) != compute_input_hash(changed)

    def test_dcf_artifact_input_hash_independent_of_generated_at(self):
        """Same run, two different timestamps -> identical input_hash."""
        inputs = _dcf_inputs()
        result = _run_dcf_default(inputs)
        art_1 = build_dcf_artifact(
            inputs=inputs, result=result, generated_at=FIXED_TIME,
            capital_structure_basis="target", discounting_convention="mid_year",
            terminal_value_method="perpetuity_growth",
        )
        art_2 = build_dcf_artifact(
            inputs=inputs, result=result, generated_at=FIXED_TIME + timedelta(days=30),
            capital_structure_basis="target", discounting_convention="mid_year",
            terminal_value_method="perpetuity_growth",
        )
        assert art_1.input_hash == art_2.input_hash
        assert art_1.generated_at != art_2.generated_at
        assert art_1.model_version == art_2.model_version

    def test_dcf_artifact_input_hash_changes_when_a_config_knob_changes(self):
        """discounting_convention is not in `inputs` but materially changes the model."""
        inputs = _dcf_inputs()
        result_mid_year = run_dcf(
            inputs, capital_structure_basis="target", discounting_convention="mid_year",
            terminal_value_method="perpetuity_growth",
        )
        result_year_end = run_dcf(
            inputs, capital_structure_basis="target", discounting_convention="year_end",
            terminal_value_method="perpetuity_growth",
        )
        art_mid = build_dcf_artifact(
            inputs=inputs, result=result_mid_year, generated_at=FIXED_TIME,
            capital_structure_basis="target", discounting_convention="mid_year",
            terminal_value_method="perpetuity_growth",
        )
        art_year_end = build_dcf_artifact(
            inputs=inputs, result=result_year_end, generated_at=FIXED_TIME,
            capital_structure_basis="target", discounting_convention="year_end",
            terminal_value_method="perpetuity_growth",
        )
        assert art_mid.input_hash != art_year_end.input_hash

    def test_comps_and_three_statement_hashes_also_move_on_one_number(self):
        target = _comps_target()
        peers = _comps_peers()
        result = run_comps(target, peers, calendarisation_policy="ltm")
        baseline = build_comps_artifact(
            target=target, peers=peers, calendarisation_policy="ltm",
            result=result, generated_at=FIXED_TIME,
        )
        bumped_target = _comps_target()
        bumped_peers = _comps_peers()
        bumped_target = TargetCompany(
            **{**bumped_target.__dict__, "total_debt": bumped_target.total_debt + 1.0}
        )
        bumped_result = run_comps(bumped_target, bumped_peers, calendarisation_policy="ltm")
        updated = build_comps_artifact(
            target=bumped_target, peers=bumped_peers, calendarisation_policy="ltm",
            result=bumped_result, generated_at=FIXED_TIME,
        )
        assert baseline.input_hash != updated.input_hash


# ---------------------------------------------------------------------------
# 2. No hidden clock
# ---------------------------------------------------------------------------


class TestNoHiddenClock:
    def test_source_never_calls_datetime_now_or_time_time(self):
        """AST-based, not substring-based: the module docstring legitimately

        *mentions* ``datetime.now()`` in prose while explaining why it is
        never called, so a naive substring grep would false-positive on the
        documentation itself. This walks the actual parsed call sites.
        """
        tree = ast.parse(inspect.getsource(artifact_module))
        offending_calls = []
        for node in ast.walk(tree):
            if not isinstance(node, ast.Call) or not isinstance(node.func, ast.Attribute):
                continue
            attr = node.func.attr
            base = node.func.value
            base_name = base.id if isinstance(base, ast.Name) else None
            if (base_name == "datetime" and attr == "now") or (base_name == "time" and attr == "time"):
                offending_calls.append(f"{base_name}.{attr}() at line {node.lineno}")
        assert offending_calls == []

    def test_missing_generated_at_is_a_typeerror_from_python_argument_binding(self):
        inputs = _dcf_inputs()
        result = _run_dcf_default(inputs)
        with pytest.raises(TypeError):
            build_dcf_artifact(  # type: ignore[call-arg]
                inputs=inputs, result=result,
                capital_structure_basis="target", discounting_convention="mid_year",
                terminal_value_method="perpetuity_growth",
            )

    def test_non_datetime_generated_at_raises_typeerror(self):
        inputs = _dcf_inputs()
        result = _run_dcf_default(inputs)
        with pytest.raises(TypeError):
            build_dcf_artifact(
                inputs=inputs, result=result, generated_at="2026-01-15",  # type: ignore[arg-type]
                capital_structure_basis="target", discounting_convention="mid_year",
                terminal_value_method="perpetuity_growth",
            )

    def test_naive_generated_at_raises_valueerror(self):
        inputs = _dcf_inputs()
        result = _run_dcf_default(inputs)
        with pytest.raises(ValueError):
            build_dcf_artifact(
                inputs=inputs, result=result, generated_at=datetime(2026, 1, 15),
                capital_structure_basis="target", discounting_convention="mid_year",
                terminal_value_method="perpetuity_growth",
            )

    def test_tz_aware_generated_at_is_accepted_and_preserved(self):
        art = _build_dcf_artifact(generated_at=FIXED_TIME)
        assert art.generated_at == FIXED_TIME


# ---------------------------------------------------------------------------
# 3. Assumption capture
# ---------------------------------------------------------------------------


class TestAssumptionCapture:
    def test_dcf_artifact_captures_both_assumptions_with_basis(self):
        art = _build_dcf_artifact()
        assert len(art.assumptions) == 2
        names = {r.assumption.name for r in art.assumptions}
        assert names == {"terminal_growth", "exit_multiple"}
        for record in art.assumptions:
            assert isinstance(record, AssumptionRecord)
            assert record.assumption.basis.strip() != ""
            assert record.path.startswith("$.terminal_value.")

    def test_dcf_artifact_assumption_values_match_the_result(self):
        art = _build_dcf_artifact()
        by_name = {r.assumption.name: r.assumption for r in art.assumptions}
        assert by_name["terminal_growth"].value == pytest.approx(0.03)
        assert by_name["exit_multiple"].value == pytest.approx(8.0)
        assert by_name["terminal_growth"].source == "analyst view"
        assert by_name["exit_multiple"].source is None

    def test_comps_artifact_has_no_assumptions(self):
        """comps.py takes no Assumption objects -- this must be an honest empty tuple."""
        _, _, _, art = _build_comps_artifact()
        assert art.assumptions == ()

    def test_three_statement_artifact_has_no_assumptions(self):
        """threestatement.py takes no Assumption objects either."""
        _, _, _, art = _build_three_statement_artifact()
        assert art.assumptions == ()


# ---------------------------------------------------------------------------
# 3b. Exclusions / omissions surfaced on the artifact itself
# ---------------------------------------------------------------------------


class TestGapExtraction:
    def test_comps_artifact_records_excluded_peer_and_omitted_components(self):
        _, _, result, art = _build_comps_artifact()
        joined = " | ".join(art.excluded)
        assert "PEER_LOSSMAKER" in joined
        assert "ev_ebitda" in joined
        assert "PEER_A" in joined  # PEER_A omits all three optional EV-bridge items
        assert "minority_interest" in joined
        assert "target" in joined.lower() or "TARGET" in joined

    def test_three_statement_artifact_records_revolver_draw(self):
        _, _, result, art = _build_three_statement_artifact()
        joined = " | ".join(art.excluded)
        assert "revolver drew" in joined

    def test_dcf_artifact_always_notes_the_non_chosen_terminal_method(self):
        art = _build_dcf_artifact()
        joined = " | ".join(art.excluded)
        assert "exit_multiple" in joined
        assert "terminal_value_method" in joined


# ---------------------------------------------------------------------------
# 4. Diff
# ---------------------------------------------------------------------------


class TestDiff:
    def test_diff_across_model_names_raises(self):
        dcf_art = _build_dcf_artifact()
        _, _, _, comps_art = _build_comps_artifact()
        with pytest.raises(ValueError):
            diff_artifacts(dcf_art, comps_art)

    def test_diff_precisely_reports_a_single_changed_input_and_its_delta(self):
        baseline_inputs = _dcf_inputs()
        baseline_art = _build_dcf_artifact(baseline_inputs)

        updated_inputs = _dcf_inputs()
        updated_inputs["total_debt"] = 250.0
        updated_art = _build_dcf_artifact(updated_inputs)

        diff = diff_artifacts(baseline_art, updated_art)
        assert diff.input_hash_changed is True
        assert diff.model_version_changed is False

        by_path = {c.path: c for c in diff.changed_inputs}
        assert "$.model_inputs.total_debt" in by_path
        change = by_path["$.model_inputs.total_debt"]
        assert change.old == pytest.approx(200.0)
        assert change.new == pytest.approx(250.0)
        assert change.delta == pytest.approx(50.0)

        # Nothing else in the raw inputs should have moved.
        other_paths = [p for p in by_path if p != "$.model_inputs.total_debt"]
        assert other_paths == []

    def test_diff_reports_only_the_outputs_actually_affected_by_total_debt(self):
        baseline_inputs = _dcf_inputs()
        baseline_art = _build_dcf_artifact(baseline_inputs)

        updated_inputs = _dcf_inputs()
        updated_inputs["total_debt"] = 250.0
        updated_art = _build_dcf_artifact(updated_inputs)

        diff = diff_artifacts(baseline_art, updated_art)
        changed_paths = {c.path for c in diff.changed_outputs}

        # total_debt flows only through the net-debt bridge to equity value /
        # value per share -- enterprise value and WACC must be untouched.
        assert "$.net_debt_bridge.total_debt" in changed_paths
        assert "$.equity_value" in changed_paths
        assert "$.value_per_share" in changed_paths
        assert "$.net_debt_bridge.equity_value" in changed_paths
        assert "$.net_debt_bridge.value_per_share" in changed_paths
        assert "$.enterprise_value" not in changed_paths
        assert "$.wacc_build.wacc" not in changed_paths

        by_path = {c.path: c for c in diff.changed_outputs}
        expected_delta = -50.0  # equity_value = EV - total_debt + ... ; +50 debt -> -50 equity
        assert by_path["$.equity_value"].delta == pytest.approx(expected_delta)
        assert by_path["$.net_debt_bridge.equity_value"].delta == pytest.approx(expected_delta)

    def test_diff_of_identical_artifacts_has_no_changes(self):
        inputs = _dcf_inputs()
        art_1 = _build_dcf_artifact(inputs, generated_at=FIXED_TIME)
        art_2 = _build_dcf_artifact(inputs, generated_at=FIXED_TIME + timedelta(hours=1))
        diff = diff_artifacts(art_1, art_2)
        assert diff.changed_inputs == ()
        assert diff.changed_outputs == ()
        assert diff.input_hash_changed is False


# ---------------------------------------------------------------------------
# 5. xlsx exports
# ---------------------------------------------------------------------------


class TestDCFWorkbookExport:
    @pytest.fixture()
    def workbook_path(self, tmp_path: Path) -> Path:
        inputs = _dcf_inputs()
        result = _run_dcf_default(inputs)
        art = build_dcf_artifact(
            inputs=inputs, result=result, generated_at=FIXED_TIME,
            capital_structure_basis="target", discounting_convention="mid_year",
            terminal_value_method="perpetuity_growth",
        )
        grid = sensitivity_grid(
            result.fcff_bridge,
            wacc_values=[0.08, 0.0905, 0.10],
            growth_values=[0.02, 0.03, 0.04],
            discounting_convention="mid_year",
            total_debt=200.0, cash_and_equivalents=50.0, minority_interest=10.0,
            preferred_equity=15.0, associate_investments=5.0, diluted_shares=100.0,
        )
        out = tmp_path / "dcf.xlsx"
        export_dcf_workbook(art, out, sensitivity_grid=grid)
        return out, art, result, grid

    def test_all_five_sheets_present(self, workbook_path):
        out, *_ = workbook_path
        wb = openpyxl.load_workbook(out)
        assert set(wb.sheetnames) == {
            "FCFF Bridge", "WACC Build", "Terminal Value", "Sensitivity",
            "Assumptions & Data Quality",
        }

    def test_fcff_bridge_cells_match_result(self, workbook_path):
        out, _art, result, _grid = workbook_path
        wb = openpyxl.load_workbook(out)
        ws = wb["FCFF Bridge"]
        assert ws["A1"].value == "Year"
        for i, year in enumerate(result.fcff_bridge):
            row = i + 2
            assert ws.cell(row=row, column=1).value == year.year
            assert ws.cell(row=row, column=2).value == pytest.approx(year.ebit)
            assert ws.cell(row=row, column=8).value == pytest.approx(year.fcff)
        total_row = len(result.fcff_bridge) + 2
        assert ws.cell(row=total_row, column=10).value == pytest.approx(result.pv_of_explicit_fcff)

    def test_wacc_sheet_cells_match_result(self, workbook_path):
        out, _art, result, _grid = workbook_path
        wb = openpyxl.load_workbook(out)
        ws = wb["WACC Build"]
        values = {row[0].value: row[1].value for row in ws.iter_rows(min_row=2)}
        assert values["wacc"] == pytest.approx(result.wacc_build.wacc)
        assert values["cost_of_equity"] == pytest.approx(result.wacc_build.cost_of_equity)
        assert values["equity_weight"] == pytest.approx(0.7)

    def test_terminal_value_sheet_cells_match_result(self, workbook_path):
        out, _art, result, _grid = workbook_path
        wb = openpyxl.load_workbook(out)
        ws = wb["Terminal Value"]
        values = {row[0].value: row[1].value for row in ws.iter_rows(min_row=2)}
        assert values["value_per_share"] == pytest.approx(result.value_per_share)
        assert values["enterprise_value"] == pytest.approx(result.enterprise_value)
        assert values["exit_multiple.basis"] == result.terminal_value.exit_multiple.basis
        assert values["growth_exceeds_gdp_ceiling"] == result.terminal_value.growth_exceeds_gdp_ceiling

    def test_sensitivity_sheet_matches_grid(self, workbook_path):
        out, _art, _result, grid = workbook_path
        wb = openpyxl.load_workbook(out)
        ws = wb["Sensitivity"]
        # Header row: label + growth columns.
        assert ws.cell(row=1, column=1).value == "wacc \\ terminal_growth"
        for col_idx, growth in enumerate(grid.columns, start=2):
            assert ws.cell(row=1, column=col_idx).value == pytest.approx(float(growth))
        for row_idx, wacc_value in enumerate(grid.index, start=2):
            assert ws.cell(row=row_idx, column=1).value == pytest.approx(float(wacc_value))
            for col_idx, growth in enumerate(grid.columns, start=2):
                expected = grid.loc[wacc_value, growth]
                cell_value = ws.cell(row=row_idx, column=col_idx).value
                if pd.isna(expected):
                    assert cell_value is None
                else:
                    assert cell_value == pytest.approx(float(expected))

    def test_assumptions_and_gaps_sheet_has_hash_and_both_assumptions(self, workbook_path):
        out, art, _result, _grid = workbook_path
        wb = openpyxl.load_workbook(out)
        ws = wb["Assumptions & Data Quality"]
        all_values = [cell.value for row in ws.iter_rows() for cell in row if cell.value is not None]
        assert art.input_hash in all_values
        assert art.model_version in all_values
        assert "terminal_growth" in all_values
        assert "exit_multiple" in all_values
        joined_text = " | ".join(str(v) for v in all_values)
        assert "Long-run nominal GDP proxy" in joined_text


class TestCompsWorkbookExport:
    @pytest.fixture()
    def workbook_path(self, tmp_path: Path):
        target, peers, result, art = _build_comps_artifact()
        out = tmp_path / "comps.xlsx"
        export_comps_workbook(art, out)
        return out, target, peers, result, art

    def test_all_four_sheets_present(self, workbook_path):
        out, *_ = workbook_path
        wb = openpyxl.load_workbook(out)
        assert set(wb.sheetnames) == {
            "Peer Detail", "Multiple Matrix", "Implied Valuation", "Assumptions & Data Quality",
        }

    def test_peer_detail_cells_match_result(self, workbook_path):
        out, _target, _peers, result, _art = workbook_path
        wb = openpyxl.load_workbook(out)
        ws = wb["Peer Detail"]
        header = [c.value for c in ws[1]]
        rows_by_name = {
            row[header.index("name")].value: row for row in ws.iter_rows(min_row=2)
        }
        peer_a = next(p for p in result.peer_multiples if p.name == "PEER_A")
        row = rows_by_name["PEER_A"]
        assert row[header.index("enterprise_value")].value == pytest.approx(peer_a.ev_bridge.enterprise_value)
        assert row[header.index("ev_ebitda")].value == pytest.approx(peer_a.ev_ebitda)
        lossmaker_row = rows_by_name["PEER_LOSSMAKER"]
        assert lossmaker_row[header.index("ev_ebitda")].value == "excluded"

    def test_multiple_matrix_includes_excluded_detail(self, workbook_path):
        out, _target, _peers, result, _art = workbook_path
        wb = openpyxl.load_workbook(out)
        ws = wb["Multiple Matrix"]
        rows = list(ws.iter_rows(values_only=True))
        header_idx = next(i for i, r in enumerate(rows) if r and r[0] == "multiple")
        by_multiple = {r[0]: r for r in rows[header_idx + 1 : header_idx + 1 + len(list(result.distributions))]}
        ev_ebitda_row = by_multiple["ev_ebitda"]
        dist = result.distributions["ev_ebitda"]
        assert ev_ebitda_row[3] == pytest.approx(dist.median)  # median column
        assert ev_ebitda_row[6] == len(dist.included)
        assert ev_ebitda_row[7] == len(dist.excluded)

        excluded_detail_rows = [r for r in rows if r and r[0] == "PEER_LOSSMAKER"]
        assert excluded_detail_rows, "excluded peer must appear in the Excluded detail block"
        assert any(r[1] == "ev_ebitda" for r in excluded_detail_rows)

    def test_implied_valuation_cells_match_result(self, workbook_path):
        out, _target, _peers, result, _art = workbook_path
        wb = openpyxl.load_workbook(out)
        ws = wb["Implied Valuation"]
        rows = {r[0]: r for r in ws.iter_rows(min_row=2, values_only=True)}
        ev_ebitda_row = rows["ev_ebitda"]
        iv = result.implied_valuations["ev_ebitda"]
        assert ev_ebitda_row[2] == pytest.approx(iv.target_metric_value)
        # Columns: multiple(0), is_ev_multiple(1), target_metric_value(2),
        # then implied_ev min/p25/median/p75/max at indices 3..7.
        assert ev_ebitda_row[5] == pytest.approx(iv.implied_ev_by_quantile["median"])

    def test_assumptions_sheet_honestly_reports_no_assumptions(self, workbook_path):
        out, _target, _peers, result, art = workbook_path
        wb = openpyxl.load_workbook(out)
        ws = wb["Assumptions & Data Quality"]
        all_values = [cell.value for row in ws.iter_rows() for cell in row if cell.value is not None]
        assert art.input_hash in all_values
        joined = " | ".join(str(v) for v in all_values)
        assert "no caller-supplied Assumption objects" in joined
        assert "PEER_LOSSMAKER" in joined
        assert "minority_interest" in joined


class TestThreeStatementWorkbookExport:
    @pytest.fixture()
    def workbook_path(self, tmp_path: Path):
        opening, drivers, result, art = _build_three_statement_artifact()
        out = tmp_path / "three_statement.xlsx"
        export_three_statement_workbook(art, out)
        return out, opening, drivers, result, art

    def test_all_four_sheets_present(self, workbook_path):
        out, *_ = workbook_path
        wb = openpyxl.load_workbook(out)
        assert set(wb.sheetnames) == {
            "Income Statement", "Cash Flow Statement", "Balance Sheet", "Assumptions & Data Quality",
        }

    def test_income_statement_cells_match_result(self, workbook_path):
        out, _opening, _drivers, result, _art = workbook_path
        wb = openpyxl.load_workbook(out)
        ws = wb["Income Statement"]
        for i, period in enumerate(result.periods):
            row = i + 2
            s = period.income_statement
            assert ws.cell(row=row, column=1).value == s.period
            assert ws.cell(row=row, column=2).value == pytest.approx(s.revenue)
            assert ws.cell(row=row, column=11).value == pytest.approx(s.net_income)

    def test_cash_flow_sheet_surfaces_revolver_draw_and_convergence(self, workbook_path):
        out, _opening, _drivers, result, _art = workbook_path
        wb = openpyxl.load_workbook(out)
        ws = wb["Cash Flow Statement"]
        header = [c.value for c in ws[1]]
        draw_col = header.index("revolver_draw") + 1
        converged_col = header.index("converged") + 1
        iterations_col = header.index("iterations") + 1
        for i, period in enumerate(result.periods):
            row = i + 2
            c = period.cash_flow_statement
            assert ws.cell(row=row, column=draw_col).value == pytest.approx(c.revolver_draw)
            assert ws.cell(row=row, column=converged_col).value == period.converged
            assert ws.cell(row=row, column=iterations_col).value == period.iterations
        # Period 1's minimum_cash (400,000) forces a draw -- must be > 0 and present.
        assert ws.cell(row=2, column=draw_col).value > 0.0

    def test_balance_sheet_rows_balance_and_match_result(self, workbook_path):
        out, _opening, _drivers, result, _art = workbook_path
        wb = openpyxl.load_workbook(out)
        ws = wb["Balance Sheet"]
        rows = list(ws.iter_rows(min_row=2, values_only=True))
        assert rows[0][0] == 0  # opening balance sheet first
        assert rows[0][1] == pytest.approx(result.opening_balance_sheet.cash)
        for row in rows:
            residual = row[-1]
            assert residual == pytest.approx(0.0, abs=1e-6)
        for i, period in enumerate(result.periods):
            row = rows[i + 1]
            assert row[1] == pytest.approx(period.balance_sheet.cash)
            assert row[4] == pytest.approx(period.balance_sheet.revolver_balance)

    def test_assumptions_sheet_reports_no_assumptions_and_revolver_note(self, workbook_path):
        out, _opening, _drivers, _result, art = workbook_path
        wb = openpyxl.load_workbook(out)
        ws = wb["Assumptions & Data Quality"]
        all_values = [cell.value for row in ws.iter_rows() for cell in row if cell.value is not None]
        assert art.input_hash in all_values
        joined = " | ".join(str(v) for v in all_values)
        assert "no caller-supplied Assumption objects" in joined
        assert "revolver drew" in joined


# ---------------------------------------------------------------------------
# 6. pptx summary export
# ---------------------------------------------------------------------------


def _all_text(prs: Presentation) -> str:
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.append(shape.text_frame.text)
    return "\n".join(chunks)


class TestSummaryDeckExport:
    def test_dcf_deck_contains_value_per_share_hash_and_assumptions(self, tmp_path):
        inputs = _dcf_inputs()
        result = _run_dcf_default(inputs)
        art = build_dcf_artifact(
            inputs=inputs, result=result, generated_at=FIXED_TIME,
            capital_structure_basis="target", discounting_convention="mid_year",
            terminal_value_method="perpetuity_growth",
        )
        grid = sensitivity_grid(
            result.fcff_bridge, wacc_values=[0.08, 0.0905, 0.10], growth_values=[0.02, 0.03],
            discounting_convention="mid_year", total_debt=200.0, cash_and_equivalents=50.0,
            minority_interest=10.0, preferred_equity=15.0, associate_investments=5.0,
            diluted_shares=100.0,
        )
        out = tmp_path / "dcf_summary.pptx"
        export_summary_deck(art, out, sensitivity_grid=grid)

        prs = Presentation(out)
        assert len(prs.slides) == 1
        text = _all_text(prs)
        assert f"{result.value_per_share:,.4f}" in text
        assert art.input_hash in text
        assert "terminal_growth" in text
        assert "exit_multiple" in text
        assert "Sensitivity grid not supplied" not in text

    def test_dcf_deck_without_grid_says_so_explicitly(self, tmp_path):
        art = _build_dcf_artifact()
        out = tmp_path / "dcf_summary_no_grid.pptx"
        export_summary_deck(art, out)
        prs = Presentation(out)
        text = _all_text(prs)
        assert "Sensitivity grid not supplied" in text

    def test_comps_deck_reports_no_assumptions_and_exclusions(self, tmp_path):
        _, _, result, art = _build_comps_artifact()
        out = tmp_path / "comps_summary.pptx"
        export_summary_deck(art, out)
        prs = Presentation(out)
        text = _all_text(prs)
        assert "No caller-supplied Assumption objects" in text
        assert "PEER_LOSSMAKER" in text
        assert art.input_hash in text

    def test_three_statement_deck_reports_final_period_and_gaps(self, tmp_path):
        _, _, result, art = _build_three_statement_artifact()
        out = tmp_path / "three_statement_summary.pptx"
        export_summary_deck(art, out)
        prs = Presentation(out)
        text = _all_text(prs)
        last = result.periods[-1]
        assert f"{last.income_statement.net_income:,.2f}" in text
        assert "not defined for a three-statement projection" in text
        assert "revolver drew" in text


# ---------------------------------------------------------------------------
# Construction guards
# ---------------------------------------------------------------------------


class TestModelArtifactConstruction:
    def test_rejects_unknown_model_name(self):
        art = _build_dcf_artifact()
        with pytest.raises(ValueError):
            ModelArtifact(
                model_name="not_a_real_model",
                model_version=art.model_version,
                schema_version=art.schema_version,
                generated_at=art.generated_at,
                input_hash=art.input_hash,
                inputs=art.inputs,
                outputs=art.outputs,
                assumptions=art.assumptions,
                result=art.result,
                excluded=art.excluded,
                warnings=art.warnings,
            )

    def test_rejects_malformed_input_hash(self):
        art = _build_dcf_artifact()
        with pytest.raises(ValueError):
            ModelArtifact(
                model_name=art.model_name,
                model_version=art.model_version,
                schema_version=art.schema_version,
                generated_at=art.generated_at,
                input_hash="not-a-hash",
                inputs=art.inputs,
                outputs=art.outputs,
                assumptions=art.assumptions,
                result=art.result,
                excluded=art.excluded,
                warnings=art.warnings,
            )

    def test_result_type_mismatch_raises(self):
        inputs = _dcf_inputs()
        with pytest.raises(TypeError):
            build_dcf_artifact(
                inputs=inputs, result="not a DCFResult", generated_at=FIXED_TIME,  # type: ignore[arg-type]
                capital_structure_basis="target", discounting_convention="mid_year",
                terminal_value_method="perpetuity_growth",
            )
